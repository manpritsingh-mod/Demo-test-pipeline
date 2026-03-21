"""
Build Director-Level Internship PPT - Manprit Singh
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# --- Config ---
IMG_DIR = r"C:\Users\dell\.gemini\antigravity\brain\d82f2c8a-ba55-4a12-bad0-e9de54167bda"
OUT_PATH = r"c:\Users\dell\Desktop\super-final-director-level -ppt\Manprit_Internship_Director_PPT.pptx"

# Colors
DARK_BG = RGBColor(0x12, 0x12, 0x2e)
TEAL = RGBColor(0x00, 0xB4, 0xD8)
BLUE = RGBColor(0x21, 0x96, 0xF3)
PURPLE = RGBColor(0x7C, 0x3A, 0xED)
GREEN = RGBColor(0x10, 0xB9, 0x81)
ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
RED = RGBColor(0xEF, 0x44, 0x44)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF1, 0xF5, 0xF9)
DARK_TEXT = RGBColor(0x1E, 0x29, 0x3B)
GRAY_TEXT = RGBColor(0x64, 0x74, 0x8B)
GOLD = RGBColor(0xF5, 0x9E, 0x0B)
LIGHT_TEAL_BG = RGBColor(0xE0, 0xF7, 0xFA)
LIGHT_GREEN_BG = RGBColor(0xE8, 0xF5, 0xE9)
LIGHT_BLUE_BG = RGBColor(0xE3, 0xF2, 0xFD)
LIGHT_PURPLE_BG = RGBColor(0xF3, 0xE8, 0xFF)
LIGHT_ORANGE_BG = RGBColor(0xFF, 0xF7, 0xED)
LIGHT_RED_BG = RGBColor(0xFE, 0xF2, 0xF2)
LIGHT_GOLD_BG = RGBColor(0xFF, 0xF8, 0xE1)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_shape(slide, left, top, width, height, fill_color=None, border_color=None, border_width=Pt(1)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color or WHITE
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=14, bold=False, color=DARK_TEXT, alignment=PP_ALIGN.LEFT, font_name='Segoe UI'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=11, color=DARK_TEXT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Segoe UI'
        p.space_after = Pt(4)
    return txBox

def add_header_bar(slide, text, slide_num):
    # Header background
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.9))
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK_BG
    bar.line.fill.background()
    # Accent line
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.88), prs.slide_width, Inches(0.04))
    accent.fill.solid()
    accent.fill.fore_color.rgb = TEAL
    accent.line.fill.background()
    # Title text
    add_text_box(slide, Inches(0.6), Inches(0.15), Inches(10), Inches(0.6), text, font_size=28, bold=True, color=WHITE)
    # Slide number
    add_text_box(slide, Inches(12), Inches(0.15), Inches(1), Inches(0.6), str(slide_num), font_size=24, bold=True, color=TEAL, alignment=PP_ALIGN.RIGHT)

def add_image_safe(slide, img_name, left, top, width, height):
    for ext in ['.png', '.jpg']:
        matches = [f for f in os.listdir(IMG_DIR) if f.startswith(img_name) and f.endswith(ext)]
        if matches:
            path = os.path.join(IMG_DIR, matches[0])
            slide.shapes.add_picture(path, left, top, width, height)
            return True
    return False

def add_notes(slide, text):
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text

def add_four_columns(slide, objective, achievements, learnings, challenges, y_start=Inches(1.2)):
    col_width = Inches(2.9)
    gap = Inches(0.2)
    x_start = Inches(0.4)
    col_height = Inches(5.8)
    
    cols_data = [
        ("🎯 OBJECTIVE", objective, GREEN, LIGHT_GREEN_BG),
        ("✅ ACHIEVEMENTS", achievements, GREEN, LIGHT_GREEN_BG),
        ("📚 LEARNINGS", learnings, PURPLE, LIGHT_PURPLE_BG),
        ("⚡ CHALLENGES", challenges, RED, LIGHT_RED_BG),
    ]
    
    for i, (header, items, hdr_color, bg_color) in enumerate(cols_data):
        x = x_start + i * (col_width + gap)
        # Column box
        box = add_shape(slide, x, y_start, col_width, col_height, fill_color=bg_color, border_color=hdr_color, border_width=Pt(2))
        # Header
        add_text_box(slide, x + Inches(0.15), y_start + Inches(0.1), col_width - Inches(0.3), Inches(0.4), header, font_size=13, bold=True, color=hdr_color)
        # Separator line
        sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Inches(0.15), y_start + Inches(0.5), col_width - Inches(0.3), Inches(0.03))
        sep.fill.solid()
        sep.fill.fore_color.rgb = hdr_color
        sep.line.fill.background()
        # Items
        add_bullet_list(slide, x + Inches(0.15), y_start + Inches(0.6), col_width - Inches(0.3), col_height - Inches(0.8), items, font_size=10)

# ============================
# SLIDE 1: TITLE
# ============================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid()
bg.fill.fore_color.rgb = DARK_BG
bg.line.fill.background()

# Accent shape top
accent_top = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.06))
accent_top.fill.solid()
accent_top.fill.fore_color.rgb = TEAL
accent_top.line.fill.background()

add_text_box(slide1, Inches(1), Inches(1.5), Inches(8), Inches(1.2), "INTERNSHIP", font_size=54, bold=True, color=WHITE)
add_text_box(slide1, Inches(1), Inches(2.5), Inches(8), Inches(0.6), "(MAY '25 – MAY '26)", font_size=26, bold=False, color=TEAL)

sep_line = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(3.3), Inches(5), Inches(0.04))
sep_line.fill.solid()
sep_line.fill.fore_color.rgb = TEAL
sep_line.line.fill.background()

add_text_box(slide1, Inches(1), Inches(3.6), Inches(8), Inches(0.5), "BUILD CI/CD SOLUTION", font_size=22, bold=True, color=LIGHT_GRAY)

# Team info box
team_box = add_shape(slide1, Inches(8.5), Inches(5), Inches(4), Inches(1.8), fill_color=RGBColor(0x1E, 0x1E, 0x3F), border_color=TEAL, border_width=Pt(2))
add_text_box(slide1, Inches(8.7), Inches(5.15), Inches(3.6), Inches(0.35), "Mentor: Viresh Hiremath", font_size=13, color=WHITE)
add_text_box(slide1, Inches(8.7), Inches(5.55), Inches(3.6), Inches(0.35), "Architect: Sunil Kumar EK", font_size=13, color=WHITE)
add_text_box(slide1, Inches(8.7), Inches(5.95), Inches(3.6), Inches(0.35), "Manager: Gautom Dutta", font_size=13, color=WHITE)

add_text_box(slide1, Inches(1), Inches(6.6), Inches(6), Inches(0.4), "PRESENTED BY: MANPRIT SINGH", font_size=12, bold=True, color=GRAY_TEXT)

# Add hero image
add_image_safe(slide1, "slide1_cicd_hero", Inches(8.5), Inches(1.2), Inches(4), Inches(3.2))

add_notes(slide1, """Good morning, I'm Manprit Singh, intern here since May 2025 under Gautom-san, working closely with Sunil-san and Viresh as my mentor.

Over the past few months, my journey has been from learning the fundamentals to building solutions that save time, bring standardization, and even exploring AI/ML for smarter decision-making.

Let me walk you through my journey.

TRANSITION: "But first, a quick bit about me..." """)

# ============================
# SLIDE 2: ABOUT ME
# ============================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_header_bar(slide2, "👤  ABOUT ME", 2)

# Name
add_text_box(slide2, Inches(5), Inches(1.2), Inches(5), Inches(0.6), "Manprit Singh Panesar", font_size=30, bold=True, color=DARK_TEXT)

# Photo placeholder
photo_box = add_shape(slide2, Inches(8.5), Inches(1.8), Inches(3.5), Inches(4.5), fill_color=LIGHT_GRAY, border_color=TEAL, border_width=Pt(3))
add_text_box(slide2, Inches(8.8), Inches(3.5), Inches(3), Inches(0.8), "📷 Paste Your Photo Here", font_size=14, color=GRAY_TEXT, alignment=PP_ALIGN.CENTER)

# Info items
info_items = [
    ("📋 ROLE", "Intern"),
    ("📍 HOMETOWN", "Jamshedpur, Jharkhand"),
    ("🎓 EDUCATION", "B.Tech (CSE), Arka Jain University"),
    ("💻 SKILLS", "Java, Spring Boot, CI/CD, Jenkins, Docker"),
    ("🎯 TECH INTEREST", "CI/CD Automation, Backend Development, Scalable Applications"),
    ("🎨 BEYOND TECH", "Art & Design, Travel, Badminton"),
]

y = Inches(2.0)
for label, value in info_items:
    add_text_box(slide2, Inches(0.8), y, Inches(7), Inches(0.3), label, font_size=12, bold=True, color=TEAL)
    add_text_box(slide2, Inches(0.8), y + Inches(0.3), Inches(7), Inches(0.3), value, font_size=13, color=DARK_TEXT)
    y += Inches(0.7)

add_notes(slide2, """Quick intro — I come from Jamshedpur, Jharkhand and completed my B.Tech in Computer Science from Arka Jain University.

My core interest is in CI/CD automation and backend development — Java, Spring Boot.

Outside work, I enjoy sketching, travel, and playing badminton.

TRANSITION: "Now, let me show you where this year has taken me..." """)

# ============================
# SLIDE 3: ROADMAP
# ============================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_header_bar(slide3, "📋  INTERNSHIP ROADMAP", 3)

# Timeline arrow
arrow = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.0), Inches(11.5), Inches(0.05))
arrow.fill.solid()
arrow.fill.fore_color.rgb = DARK_TEXT
arrow.line.fill.background()

quarters = [
    ("Q1 FY '25", "🧱 Foundation & POC 1", ["Jenkins & Docker fundamentals", "Groovy & Python scripting", "POC 1: Repository Maintenance"], LIGHT_GREEN_BG, GREEN),
    ("Q2 FY '25", "🔗 Unified CI Core (POC 2)", ["Shared Library architecture", "Java & Python pipelines", "POC 2: Unified CI"], LIGHT_GREEN_BG, GREEN),
    ("Q3 FY '25", "⚛️ React + AI/ML", ["React/React-Native support", "Jenkins MCP exploration", "AI-assisted CI research"], LIGHT_GREEN_BG, GREEN),
    ("Q4 FY '26", "📱 Mobile CI/CD + AI/ML", ["Android CI/CD pipeline", "Nexus Repo deployment", "AI/ML Integration", "🏆 Hackathon: Self-Healing CI/CD"], LIGHT_BLUE_BG, BLUE),
]

x_start = Inches(0.5)
q_width = Inches(2.9)
q_gap = Inches(0.25)

for i, (label, title, items, bg_color, border_color) in enumerate(quarters):
    x = x_start + i * (q_width + q_gap)
    # Quarter badge
    badge = add_shape(slide3, x + Inches(0.7), Inches(1.4), Inches(1.3), Inches(0.35), fill_color=WHITE, border_color=DARK_TEXT, border_width=Pt(2))
    add_text_box(slide3, x + Inches(0.75), Inches(1.4), Inches(1.2), Inches(0.35), label, font_size=10, bold=True, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)
    # Content box
    box = add_shape(slide3, x, Inches(2.5), q_width, Inches(4.2), fill_color=bg_color, border_color=border_color, border_width=Pt(2))
    add_text_box(slide3, x + Inches(0.15), Inches(2.65), q_width - Inches(0.3), Inches(0.4), title, font_size=13, bold=True, color=DARK_TEXT)
    add_bullet_list(slide3, x + Inches(0.15), Inches(3.15), q_width - Inches(0.3), Inches(3.3), items, font_size=10)

add_notes(slide3, """Here's a high-level view of my one-year journey, structured across four quarters.

In Q1, I focused on learning the fundamentals and delivered my first proof of concept.
In Q2, I moved towards solving a larger problem — standardizing CI/CD through UnifiedCI.
Q3 was about extending UnifiedCI to React projects and exploring how AI/ML can help in CI/CD.
And in Q4, I focused on innovation — ML-based node selection, mobile CI/CD, and I also participated in a hackathon on Self-Healing CI/CD with Agentic AI.

So the journey has been — learn, build, extend, and now innovate.

TRANSITION: "Let me start from the beginning — Q1..." """)

# ============================
# SLIDE 4: Q1 FOUNDATION & POC 1
# ============================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_header_bar(slide4, "Q1 FY '25  |  FOUNDATION & POC 1", 4)

add_four_columns(slide4,
    objective=["Build foundational CI/CD skills", "Understand Jenkins pipeline development", "Automate repository maintenance"],
    achievements=["Completed onboarding & KT", "Built shared-library for email/slack notify", "Delivered POC 1: Automated PR cleanup", "Presented MVP to the mentor"],
    learnings=["TECH STACK:", "Jenkins & Docker", "Groovy & Python", "Shell scripting", "", "SOFT SKILLS:", "Adaptability to corporate work culture", "Effective communication in meetings", "Asking the right questions", "Time management"],
    challenges=["Access for Slack/email token key", "GitHub Access", "Access for Docker", "Linux Desktop setup", "SWAN2 network access"]
)

add_image_safe(slide4, "slide4_automation", Inches(9.5), Inches(5.5), Inches(3.5), Inches(1.8))

add_notes(slide4, """When I joined, my mentor assigned me a problem to work on.

We had pull requests — which are basically proposed code changes waiting for review — sitting idle across repositories. Old code. Dead branches. Nobody cleaning them up.

So I built an automation where every day, a Jenkins job finds PRs inactive for over 30 days, labels them stale, closes them, and deletes the dead branches. Zero human involvement.

The biggest challenge wasn't technical. It was getting access — Linux desktop, GitHub, Docker, network approvals.

This taught me — technical knowledge is only one part of the work. Communication, follow-ups, and patience matter equally.

TRANSITION: "Let me show you how this actually works..." """)

# ============================
# SLIDE 5: Q1 ARCHITECTURE & DEMO
# ============================
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_header_bar(slide5, "Q1 FY '25  |  POC 1: ARCHITECTURE, DOC & DEMO", 5)

# Architecture flow
add_text_box(slide5, Inches(0.5), Inches(1.2), Inches(5.5), Inches(0.4), "AUTOMATED PR MANAGEMENT FLOW", font_size=16, bold=True, color=DARK_TEXT)

steps = ["1️⃣  Initialize (Daily Job Trigger)", "2️⃣  Analyze (Identify PRs > 30 Days)", "3️⃣  Label (Tag 'Stale-PR')", "4️⃣  Cleanup (Close PR & delete branch)"]
y = Inches(1.8)
for step in steps:
    step_box = add_shape(slide5, Inches(0.8), y, Inches(4.5), Inches(0.55), fill_color=LIGHT_TEAL_BG, border_color=TEAL, border_width=Pt(2))
    add_text_box(slide5, Inches(1.0), y + Inches(0.08), Inches(4), Inches(0.4), step, font_size=12, bold=True, color=DARK_TEXT)
    y += Inches(0.8)
    if step != steps[-1]:
        add_text_box(slide5, Inches(2.8), y - Inches(0.28), Inches(0.5), Inches(0.3), "⬇", font_size=16, color=TEAL, alignment=PP_ALIGN.CENTER)

# Demo placeholder
demo_box = add_shape(slide5, Inches(6.2), Inches(1.2), Inches(6.5), Inches(4.5), fill_color=LIGHT_GRAY, border_color=BLUE, border_width=Pt(2))
add_text_box(slide5, Inches(6.4), Inches(1.3), Inches(6), Inches(0.4), "🎬  POC 1: LIVE JENKINS DEMO", font_size=14, bold=True, color=BLUE)
add_text_box(slide5, Inches(7), Inches(3), Inches(5), Inches(1), "📷 Paste Jenkins Stage View\nScreenshot Here", font_size=14, color=GRAY_TEXT, alignment=PP_ALIGN.CENTER)

# Documentation box
doc_box = add_shape(slide5, Inches(0.5), Inches(5.5), Inches(6), Inches(1.5), fill_color=LIGHT_BLUE_BG, border_color=BLUE, border_width=Pt(2))
add_text_box(slide5, Inches(0.7), Inches(5.6), Inches(5.5), Inches(0.3), "📄 DOCUMENTATION", font_size=12, bold=True, color=BLUE)
add_bullet_list(slide5, Inches(0.7), Inches(5.95), Inches(5.5), Inches(0.9),
    ["GITHUB LINK - https://github.sie.sony.com/7000041870/auto_PR_cleanUp.git", "PPT - Phase 1 (Q1 2025) Presentation"],
    font_size=9, color=DARK_TEXT)

add_notes(slide5, """Here's the system I just described.

It works in four simple steps — all automatic. Every day, a Jenkins job triggers on its own, scans every open pull request, finds the ones inactive for over 30 days, labels them stale, then closes them and deletes the dead branches.

And here is the video demo — the actual Jenkins pipeline running. Real runs. Each stage finishes in seconds.

Something that was manual and easily forgotten now runs like clockwork, every single day.

TRANSITION: "That was a good start. But in Q2, I noticed a much bigger problem..." """)

# ============================
# SLIDE 6: Q2 UNIFIED-CI CORE
# ============================
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
add_header_bar(slide6, "Q2 FY '25  |  UNIFIED-CI CORE", 6)

add_four_columns(slide6,
    objective=["Build unified Jenkins Shared Library", "Support Java & Python", "Standardize CI/CD across org"],
    achievements=["Designed & delivered Unified-CI", "Auto-Language detection", "Java (Maven/Gradle), Python support", "95% reduction in pipeline code"],
    learnings=["TECH STACK:", "Shared Library structure", "JUnit, Lint tools", "Nexus Repo integrate", "", "SOFT SKILLS:", "Collaboration with mentor", "Active listening", "Effective communication", "Receiving & implementing feedback"],
    challenges=["GitHub access", "Linux PC setup for Docker & Jenkins", "SWAN2 network access", "Tool installation access", "Disk space issues"]
)

add_image_safe(slide6, "slide6_unified_teams", Inches(9.5), Inches(5.5), Inches(3.5), Inches(1.8))

add_notes(slide6, """After a few months, we noticed a pattern. Every team was building their CI/CD pipeline from scratch. Same goal — build, test, report — but everyone doing it differently.

Imagine ten teams drawing their own maps to the same destination. So I built UnifiedCI — a shared library. Think of it as a common template that any team can reuse. A team just says "Java" or "Python" in one config file, and everything is ready.

Result? Setup went from days to hours. Pipeline code dropped. Every team now follows the same standard.

I also learned to listen before building — understanding user needs is as important as building the solution.

TRANSITION: "Let me show you the architecture..." """)

# ============================
# SLIDE 7: Q2 ARCHITECTURE & DEMO
# ============================
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
add_header_bar(slide7, "Q2 FY '25  |  ARCHITECTURE, DOC & DEMO", 7)

# Architecture
add_text_box(slide7, Inches(0.5), Inches(1.2), Inches(5.5), Inches(0.4), "UNIFIED-CI SHARED LIBRARY ARCHITECTURE", font_size=14, bold=True, color=DARK_TEXT)

# Simple architecture flow
flow_items = [
    ("Jenkinsfile", Inches(0.8), LIGHT_GRAY),
    ("Shared Library\n(Unified-CI)", Inches(2.8), LIGHT_TEAL_BG),
]
for label, x, color in flow_items:
    box = add_shape(slide7, x, Inches(2.2), Inches(1.8), Inches(1.0), fill_color=color, border_color=TEAL, border_width=Pt(2))
    add_text_box(slide7, x + Inches(0.1), Inches(2.3), Inches(1.6), Inches(0.8), label, font_size=10, bold=True, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)

# Arrow
add_text_box(slide7, Inches(2.55), Inches(2.5), Inches(0.4), Inches(0.4), "➡", font_size=16, color=TEAL, alignment=PP_ALIGN.CENTER)

# Templates
templates = [("Python\nTemplate 🐍", Inches(5.2), Inches(1.8)), ("Java\nTemplate ☕", Inches(5.2), Inches(3.0))]
for label, x, y in templates:
    box = add_shape(slide7, x, y, Inches(1.6), Inches(0.9), fill_color=LIGHT_GREEN_BG, border_color=GREEN, border_width=Pt(2))
    add_text_box(slide7, x + Inches(0.1), y + Inches(0.1), Inches(1.4), Inches(0.7), label, font_size=10, bold=True, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)

# Arrows to templates
add_text_box(slide7, Inches(4.5), Inches(2.0), Inches(0.8), Inches(0.4), "➡", font_size=14, color=TEAL)
add_text_box(slide7, Inches(4.5), Inches(3.2), Inches(0.8), Inches(0.4), "➡", font_size=14, color=TEAL)

# Demo placeholder
demo_box = add_shape(slide7, Inches(7.2), Inches(1.2), Inches(5.5), Inches(4.5), fill_color=LIGHT_GRAY, border_color=BLUE, border_width=Pt(2))
add_text_box(slide7, Inches(7.4), Inches(1.3), Inches(5), Inches(0.4), "🎬  POC 2: UNIFIEDCI SHARED LIBRARY", font_size=13, bold=True, color=BLUE)
add_text_box(slide7, Inches(8), Inches(3), Inches(4), Inches(1), "📷 Paste Jenkins Build\nScreenshot Here", font_size=14, color=GRAY_TEXT, alignment=PP_ALIGN.CENTER)

# Documentation
doc_box = add_shape(slide7, Inches(0.5), Inches(5.5), Inches(6.2), Inches(1.5), fill_color=LIGHT_BLUE_BG, border_color=BLUE, border_width=Pt(2))
add_text_box(slide7, Inches(0.7), Inches(5.6), Inches(5.8), Inches(0.3), "📄 DOCUMENTATION", font_size=12, bold=True, color=BLUE)
add_bullet_list(slide7, Inches(0.7), Inches(5.95), Inches(5.8), Inches(0.9),
    ["GITHUB - https://github.sie.sony.com/7000041870/UnifiedCI.git", "DETAILED DOC - Detail Report (Q2+Q3)", "PPT - Phase 2 (Q2) Presentation"],
    font_size=9, color=DARK_TEXT)

add_notes(slide7, """Here's the architecture — simple by design.

A developer writes a small Jenkinsfile — which is just a small configuration file that tells Jenkins what to do — it connects to our shared library, the library picks the right template — Java or Python — and handles everything automatically.

Here is the demo video — a real build. Java project, completed successfully in under two minutes.

The key point — teams don't need to understand the complexity underneath. A few lines of configuration, and it just works.

TRANSITION: "Q2 was about standardization. Q3 is where things got more interesting..." """)

# ============================
# SLIDE 8: Q3 REACT + AI/ML
# ============================
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
add_header_bar(slide8, "Q3 FY '25  |  REACT + AI/ML EXPLORATION", 8)

add_four_columns(slide8,
    objective=["Extend UnifiedCI for React", "Explore AI/ML for CI automation", "Research on Jenkins MCP"],
    achievements=["React support added", "Jenkins MCP exploration", "Documented AI use-case", "Proof of concept validated"],
    learnings=["TECH STACK:", "React (npm, Jest, ESLint)", "MCP usage", "LLM use-case in CI/CD", "", "SOFT SKILLS:", "Problem-solving mindset", "Handling deadlines under pressure", "Prioritization skills", "Research & experimentation"],
    challenges=["Linux Desktop space issue", "New Desktop provisioning", "Access to GitHub Copilot", "AI tool access request"]
)

add_image_safe(slide8, "slide8_ai_brain", Inches(9.5), Inches(5.5), Inches(3.5), Inches(1.8))

add_notes(slide8, """In Q3, I extended UnifiedCI to support React projects. Now we have React, Java, Python — three technologies, one system.

But the bigger story of Q3 is: Can our CI system not just run things, but actually understand what went wrong?

When a build fails, someone reads through hundreds of lines of logs. What if the system could summarize the failure, suggest the cause, and recommend a fix?

That's where I connected GitHub Copilot to Jenkins through MCP — Model Context Protocol. Think of MCP as a bridge — it lets an AI model talk to Jenkins and provide intelligent feedback.

It's not just automation anymore. It's intelligent automation.

TRANSITION: "Let me show you the architecture..." """)

# ============================
# SLIDE 9: Q3 ARCHITECTURE & DEMO
# ============================
slide9 = prs.slides.add_slide(prs.slide_layouts[6])
add_header_bar(slide9, "Q3 FY '25  |  ARCHITECTURE, DOC & DEMO", 9)

# React Pipeline Architecture
add_text_box(slide9, Inches(0.5), Inches(1.2), Inches(5.5), Inches(0.35), "REACT PIPELINE ARCHITECTURE", font_size=14, bold=True, color=DARK_TEXT)

react_flow = [("Jenkinsfile", Inches(0.8)), ("Shared Library\n(Unified-CI)", Inches(2.8)), ("React\nTemplate ⚛️", Inches(4.8))]
for label, x in react_flow:
    box = add_shape(slide9, x, Inches(1.8), Inches(1.7), Inches(0.9), fill_color=LIGHT_TEAL_BG, border_color=TEAL, border_width=Pt(2))
    add_text_box(slide9, x + Inches(0.1), Inches(1.85), Inches(1.5), Inches(0.8), label, font_size=9, bold=True, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)

add_text_box(slide9, Inches(2.45), Inches(2.0), Inches(0.4), Inches(0.4), "➡", font_size=14, color=TEAL)
add_text_box(slide9, Inches(4.45), Inches(2.0), Inches(0.4), Inches(0.4), "➡", font_size=14, color=TEAL)

# MCP Architecture
add_text_box(slide9, Inches(0.5), Inches(3.2), Inches(5.5), Inches(0.35), "JENKINS MCP ARCHITECTURE", font_size=14, bold=True, color=DARK_TEXT)

mcp_flow = [("🧠 LLM", Inches(0.8)), ("🔌 MCP", Inches(2.8)), ("⚙️ Jenkins", Inches(4.8))]
for label, x in mcp_flow:
    box = add_shape(slide9, x, Inches(3.8), Inches(1.7), Inches(0.9), fill_color=LIGHT_PURPLE_BG, border_color=PURPLE, border_width=Pt(2))
    add_text_box(slide9, x + Inches(0.1), Inches(3.9), Inches(1.5), Inches(0.7), label, font_size=12, bold=True, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)

add_text_box(slide9, Inches(2.45), Inches(4.0), Inches(0.4), Inches(0.4), "↔", font_size=14, color=PURPLE)
add_text_box(slide9, Inches(4.45), Inches(4.0), Inches(0.4), Inches(0.4), "↔", font_size=14, color=PURPLE)

# Demo placeholder
demo_box = add_shape(slide9, Inches(7), Inches(1.2), Inches(5.8), Inches(4.2), fill_color=LIGHT_GRAY, border_color=BLUE, border_width=Pt(2))
add_text_box(slide9, Inches(7.2), Inches(1.3), Inches(5.4), Inches(0.4), "🎬  REACT PIPELINE DEMO", font_size=13, bold=True, color=BLUE)
add_text_box(slide9, Inches(8), Inches(3), Inches(4), Inches(1), "📷 Paste React Pipeline\nJenkins Screenshot Here", font_size=14, color=GRAY_TEXT, alignment=PP_ALIGN.CENTER)

# Documentation
doc_box = add_shape(slide9, Inches(0.5), Inches(5.5), Inches(6.2), Inches(1.5), fill_color=LIGHT_BLUE_BG, border_color=BLUE, border_width=Pt(2))
add_text_box(slide9, Inches(0.7), Inches(5.6), Inches(5.8), Inches(0.3), "📄 DOCUMENTATION", font_size=12, bold=True, color=BLUE)
add_bullet_list(slide9, Inches(0.7), Inches(5.95), Inches(5.8), Inches(0.9),
    ["GITHUB - https://github.sie.sony.com/7000041870/UnifiedCI.git", "PPT FOR MCP - DevOps(MCP)", "PPT FOR REACT - Phase 2 (Q2) Presentation"],
    font_size=9, color=DARK_TEXT)

add_notes(slide9, """Two things here.

On the left — the React pipeline. Same pattern. Jenkinsfile connects to the shared library, now with a React template. Standardized.

On the right — the AI piece. An AI model — called an LLM — talks to MCP, and MCP connects to Jenkins. This lets the AI pull build data and analyze logs. Still in exploration, but early results are promising.

By end of Q3 — three tech stacks covered, and AI entering our CI/CD process.

TRANSITION: "And that brings us to the current quarter..." """)

# ============================
# SLIDE 10: Q4 MOBILE CI/CD + AI/ML
# ============================
slide10 = prs.slides.add_slide(prs.slide_layouts[6])
add_header_bar(slide10, "Q4 FY '26  |  MOBILE CI/CD + AI/ML", 10)

add_four_columns(slide10,
    objective=["ML-powered AWS EC2 node selection", "Build Android CI/CD", "Enable Nexus Repo deploy"],
    achievements=["Docker environment ready", "Build & test pipeline working", "🔄 Nexus Repo deploy in progress", "Synthetic dataset generated", "🔄 Prediction model in progress"],
    learnings=["TECH STACK:", "Android CI/CD + Fastlane", "APK/AAB generation", "ML dataset preparation", "", "SOFT SKILLS:", "Time Management", "Research & experimentation", "Presentation & demo skills", "Working under hackathon deadlines"],
    challenges=["Internal Dataset needed for ML", "AWS access for node selection", "Play Store credentials needed", "Physical device for E2E test"]
)

# Hackathon callout box
hack_box = add_shape(slide10, Inches(3.55), Inches(5.8), Inches(2.7), Inches(1.1), fill_color=LIGHT_GOLD_BG, border_color=GOLD, border_width=Pt(2))
add_text_box(slide10, Inches(3.65), Inches(5.85), Inches(2.5), Inches(0.25), "🏆 HACKATHON PARTICIPATION", font_size=9, bold=True, color=RGBColor(0x7D, 0x66, 0x08))
add_text_box(slide10, Inches(3.65), Inches(6.1), Inches(2.5), Inches(0.25), "Self-Healing CI/CD with Agentic AI", font_size=8, color=RGBColor(0x8A, 0x7A, 0x2E))
add_text_box(slide10, Inches(3.65), Inches(6.35), Inches(2.5), Inches(0.25), "✅ Submitted on time", font_size=8, bold=True, color=GREEN)

add_image_safe(slide10, "slide10_mobile_ml", Inches(9.5), Inches(5.5), Inches(3.5), Inches(1.8))

add_notes(slide10, """In the current quarter, I'm working on two things.

First — ML-powered AWS node selection. Teams pick AWS machines — these are computers in the cloud — by guessing. Too big = wasted money. Too small = slow builds. What if we predicted the right machine automatically?

Right resource. Right time. No waste. Even a small optimization repeated across thousands of builds creates real savings.

Second — mobile CI/CD. Developer pushes code, everything happens automatically. Build. Test. Package. Deploy.

I also participated in a hackathon — Self-Healing CI/CD with Agentic AI. The idea: when a build fails, the system automatically understands the failure and tries to fix itself. Submitted on time.

TRANSITION: "Let me show you both architectures..." """)

# ============================
# SLIDE 11: Q4 ARCHITECTURE & DOCS
# ============================
slide11 = prs.slides.add_slide(prs.slide_layouts[6])
add_header_bar(slide11, "Q4 FY '26  |  ARCHITECTURE & DOCUMENTATION", 11)

# Mobile CI/CD Architecture
add_text_box(slide11, Inches(0.5), Inches(1.2), Inches(6), Inches(0.35), "MOBILE CI/CD ARCHITECTURE", font_size=14, bold=True, color=DARK_TEXT)

mobile_flow = [
    ("👨‍💻\nDeveloper", Inches(0.5)),
    ("📦 Source\nGitHub Org", Inches(2.0)),
    ("⚙️ Jenkins\nController", Inches(3.5)),
    ("🐳 Docker\nBuild & Test", Inches(5.0)),
    ("📱 Android\nEmulator", Inches(6.5)),
]
for label, x in mobile_flow:
    box = add_shape(slide11, x, Inches(1.8), Inches(1.3), Inches(0.9), fill_color=LIGHT_TEAL_BG, border_color=TEAL, border_width=Pt(2))
    add_text_box(slide11, x + Inches(0.05), Inches(1.85), Inches(1.2), Inches(0.8), label, font_size=8, bold=True, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)

for x in [Inches(1.75), Inches(3.25), Inches(4.75), Inches(6.25)]:
    add_text_box(slide11, x, Inches(2.0), Inches(0.3), Inches(0.3), "→", font_size=14, color=TEAL)

# ML Node Selection
add_text_box(slide11, Inches(0.5), Inches(3.3), Inches(6), Inches(0.35), "ML-POWERED AWS EC2 NODE SELECTION", font_size=14, bold=True, color=DARK_TEXT)

ml_steps = [("👨‍💻 Developer\nCode Change", Inches(0.5)), ("🚀 Build\nStart", Inches(2.5)), ("🧠 Smart Decision\nMemory & CPU", Inches(4.5))]
for label, x in ml_steps:
    box = add_shape(slide11, x, Inches(3.9), Inches(1.7), Inches(0.9), fill_color=LIGHT_PURPLE_BG, border_color=PURPLE, border_width=Pt(2))
    add_text_box(slide11, x + Inches(0.05), Inches(3.95), Inches(1.6), Inches(0.8), label, font_size=9, bold=True, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)

add_text_box(slide11, Inches(2.15), Inches(4.1), Inches(0.4), Inches(0.4), "→", font_size=14, color=PURPLE)
add_text_box(slide11, Inches(4.15), Inches(4.1), Inches(0.4), Inches(0.4), "→", font_size=14, color=PURPLE)

# Node options
nodes = [("T3a Small\n2 GB", Inches(7)), ("T3a Large\n8 GB", Inches(8.8)), ("T3a XLarge\n16 GB", Inches(10.6)), ("T3a 2XLarge\n32 GB", Inches(10.6))]
node_data = [("T3a Small\n2 GB", Inches(6.8)), ("T3a Large\n8 GB", Inches(8.3)), ("T3a XLarge\n16 GB", Inches(9.8)), ("T3a 2XL\n32 GB", Inches(11.3))]
for label, x in node_data:
    box = add_shape(slide11, x, Inches(3.9), Inches(1.3), Inches(0.9), fill_color=LIGHT_ORANGE_BG, border_color=ORANGE, border_width=Pt(2))
    add_text_box(slide11, x + Inches(0.05), Inches(3.95), Inches(1.2), Inches(0.8), label, font_size=9, bold=True, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)

# Documentation
doc_box = add_shape(slide11, Inches(0.5), Inches(5.5), Inches(6.2), Inches(1.5), fill_color=LIGHT_BLUE_BG, border_color=BLUE, border_width=Pt(2))
add_text_box(slide11, Inches(0.7), Inches(5.6), Inches(5.8), Inches(0.3), "📄 DOCUMENTATION", font_size=12, bold=True, color=BLUE)
add_bullet_list(slide11, Inches(0.7), Inches(5.95), Inches(5.8), Inches(0.9),
    ["GITHUB - https://github.sie.sony.com/7000041870/mobile-ci-cd.git",
     "GITHUB - https://github.sie.sony.com/7000041870/ML_AWS-EC2-Node-Selector.git",
     "PPT FOR ML - ML-Powered AWS EC2 Node",
     "PPT FOR MOBILE CI/CD"],
    font_size=9, color=DARK_TEXT)

add_notes(slide11, """On the left — Mobile CI/CD. Developer pushes code, Jenkins spins up a Docker container — think of it as a clean, isolated environment, like a fresh computer every time — builds the app, runs tests, and deploys to an Android emulator.

On the right — ML Node Selection. Build starts, system checks the project's history, predicts the right machine — from 2 GB small to 32 GB extra-large. The right fit, automatically chosen.

What excites me — one expands into mobile, the other brings intelligence into infrastructure.

TRANSITION: "So, let me bring it all together..." """)

# ============================
# SLIDE 12: THANK YOU
# ============================
slide12 = prs.slides.add_slide(prs.slide_layouts[6])
bg12 = slide12.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg12.fill.solid()
bg12.fill.fore_color.rgb = DARK_BG
bg12.line.fill.background()

accent12 = slide12.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(3.5), prs.slide_width, Inches(0.04))
accent12.fill.solid()
accent12.fill.fore_color.rgb = TEAL
accent12.line.fill.background()

add_text_box(slide12, Inches(0), Inches(2.5), prs.slide_width, Inches(1), "THANK YOU", font_size=52, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_text_box(slide12, Inches(0), Inches(3.8), prs.slide_width, Inches(0.6), "Happy to take any questions 🙌", font_size=22, color=TEAL, alignment=PP_ALIGN.CENTER)

# Journey recap
add_text_box(slide12, Inches(0), Inches(5.0), prs.slide_width, Inches(0.5), "🧹 → 🔗 → ⚛️🧠 → 📱💰", font_size=26, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
add_text_box(slide12, Inches(0), Inches(5.5), prs.slide_width, Inches(0.4), "Learn  →  Build  →  Extend  →  Innovate", font_size=16, color=GRAY_TEXT, alignment=PP_ALIGN.CENTER)

add_image_safe(slide12, "slide12_thankyou", Inches(9), Inches(0.5), Inches(3.5), Inches(3))

add_notes(slide12, """When I started, my focus was learning different tools — how does Jenkins work? What is Docker?

Now, my focus is solving problems — how do we stop repeated work? How do we spend less on cloud machines?

From cleanup automation... to a shared library for the org... to AI-assisted CI/CD... to intelligent infrastructure and mobile deployment.

From learning tools — to solving problems.

I want to keep building systems that are not just automated — but intelligent, scalable, and impactful.

Thank you. Happy to take your questions.""")

# ============================
# SAVE
# ============================
prs.save(OUT_PATH)
print(f"✅ PPT saved to: {OUT_PATH}")
print(f"📊 Total slides: {len(prs.slides)}")
